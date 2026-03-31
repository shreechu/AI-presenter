"""
Slide deck management — loading, navigation, pause/resume, relevance search.

Supports two modes:
  1. Load from a real *.pptx* file via ``python-pptx`` (extracts titles + speaker notes).
  2. Use a built-in demo deck for local testing.
"""

from __future__ import annotations

import asyncio
import logging
import math
import re
from collections import Counter
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Coroutine, List, Optional

logger = logging.getLogger(__name__)


# ── Data ──────────────────────────────────────────────────────────────────────

@dataclass
class Slide:
    index: int
    title: str
    speaker_notes: str


# ── Controller ────────────────────────────────────────────────────────────────

class SlideController:
    """Async-safe slide state machine with navigation, pause/resume, and search."""

    def __init__(self, slides: List[Slide]) -> None:
        if not slides:
            raise ValueError("At least one slide is required")
        self._slides = slides
        self._current_idx = 0
        self._paused = False
        self._lock = asyncio.Lock()
        self._on_change: Optional[Callable[[Slide], Coroutine]] = None

        # Pre-compute IDF for relevance scoring
        self._idf = _compute_idf([f"{s.title} {s.speaker_notes}" for s in slides])

    # ── Event hooks ───────────────────────────────────────────────────────────

    def on_slide_change(self, callback: Callable[[Slide], Coroutine]) -> None:
        """Register an async callback fired whenever the current slide changes."""
        self._on_change = callback

    async def _fire_change(self, slide: Slide) -> None:
        if self._on_change:
            try:
                await self._on_change(slide)
            except Exception:
                logger.exception("Slide-change callback failed")

    # ── State queries ─────────────────────────────────────────────────────────

    @property
    def total_slides(self) -> int:
        return len(self._slides)

    async def current_slide(self) -> Slide:
        async with self._lock:
            return self._slides[self._current_idx]

    async def progress(self) -> tuple[int, int]:
        """Return ``(current_index, total)``."""
        async with self._lock:
            return self._current_idx, len(self._slides)

    async def is_paused(self) -> bool:
        async with self._lock:
            return self._paused

    # ── Navigation ────────────────────────────────────────────────────────────

    async def pause(self) -> None:
        async with self._lock:
            self._paused = True
        logger.info("Slides paused")

    async def resume(self) -> None:
        async with self._lock:
            self._paused = False
        logger.info("Slides resumed")

    async def advance(self) -> Optional[Slide]:
        async with self._lock:
            if self._current_idx + 1 >= len(self._slides):
                return None
            self._current_idx += 1
            slide = self._slides[self._current_idx]
        await self._fire_change(slide)
        return slide

    async def go_back(self) -> Optional[Slide]:
        async with self._lock:
            if self._current_idx <= 0:
                return None
            self._current_idx -= 1
            slide = self._slides[self._current_idx]
        await self._fire_change(slide)
        return slide

    async def jump_to(self, slide_index: int) -> Slide:
        async with self._lock:
            if slide_index < 0 or slide_index >= len(self._slides):
                raise IndexError(f"Slide index {slide_index} out of range [0, {len(self._slides)})")
            self._current_idx = slide_index
            slide = self._slides[self._current_idx]
        await self._fire_change(slide)
        return slide

    async def has_next(self) -> bool:
        async with self._lock:
            return self._current_idx + 1 < len(self._slides)

    # ── Relevance search (TF-IDF) ────────────────────────────────────────────

    async def find_relevant_slide(self, question: str) -> Optional[Slide]:
        """Return the most relevant slide for *question* using TF-IDF cosine similarity."""
        q_tokens = _tokenize(question)
        if not q_tokens:
            return None

        best_slide: Optional[Slide] = None
        best_score = 0.0
        async with self._lock:
            for slide in self._slides:
                doc_tokens = _tokenize(f"{slide.title} {slide.speaker_notes}")
                score = _tfidf_similarity(q_tokens, doc_tokens, self._idf)
                if score > best_score:
                    best_score = score
                    best_slide = slide

        if best_score <= 0:
            return None
        logger.debug("Best slide match: idx=%s score=%.3f", best_slide.index if best_slide else "?", best_score)
        return best_slide


# ── Loader: PowerPoint ────────────────────────────────────────────────────────

def load_slides_from_pptx(path: str | Path) -> List[Slide]:
    """Extract slides with titles and speaker notes from a *.pptx* file."""
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError:
        raise ImportError(
            "python-pptx is required to load slides from PPTX files. "
            "Install it with: pip install python-pptx"
        )

    prs = Presentation(str(path))
    slides: List[Slide] = []
    for idx, pptx_slide in enumerate(prs.slides):
        # Title
        title = ""
        if pptx_slide.shapes.title:
            title = pptx_slide.shapes.title.text.strip()
        if not title:
            title = f"Slide {idx + 1}"

        # Speaker notes
        notes = ""
        if pptx_slide.has_notes_slide:
            notes_frame = pptx_slide.notes_slide.notes_text_frame
            notes = notes_frame.text.strip() if notes_frame else ""

        slides.append(Slide(index=idx, title=title, speaker_notes=notes))
        logger.info("Loaded slide %d: %s (%d chars of notes)", idx, title, len(notes))

    if not slides:
        raise ValueError(f"No slides found in {path}")
    return slides


# ── Loader: demo deck ────────────────────────────────────────────────────────

def demo_slides() -> List[Slide]:
    """Built-in demo deck for testing without a PPTX file."""
    return [
        Slide(
            index=0,
            title="Vision and Problem",
            speaker_notes=(
                "Welcome everyone. We are building an AI-powered Teams presenter bot that can "
                "narrate slides, monitor audience signals, and adapt in real time. "
                "The problem we solve is that live presentations are hard to scale and often "
                "lose audience engagement when the speaker cannot react to questions."
            ),
        ),
        Slide(
            index=1,
            title="Architecture",
            speaker_notes=(
                "The architecture is event-driven and modular, with independent components for "
                "slide control, speech detection, question classification, and Teams integration. "
                "Each module communicates through async queues, making it easy to replace any "
                "component — for example, swapping local TTS for Azure Neural Voices."
            ),
        ),
        Slide(
            index=2,
            title="Audience Interaction",
            speaker_notes=(
                "When an audience member speaks or types a question, the bot pauses the "
                "presentation, classifies the input, finds the most relevant slide, and "
                "generates a concise answer using the slide notes plus optional LLM context."
            ),
        ),
        Slide(
            index=3,
            title="Scalability and Deployment",
            speaker_notes=(
                "The bot can be containerised and deployed to Azure App Service or AKS. "
                "Multiple instances can handle parallel meetings. All configuration is "
                "driven by environment variables for easy CI/CD integration."
            ),
        ),
        Slide(
            index=4,
            title="Safety and Reliability",
            speaker_notes=(
                "The system logs every important event and handles API or playback errors "
                "gracefully so the presentation can continue without crashing. "
                "All external calls include timeout guards and structured error handling."
            ),
        ),
        Slide(
            index=5,
            title="Summary and Q&A",
            speaker_notes=(
                "Thank you for attending. To summarise: we demonstrated an end-to-end "
                "AI presenter that advances slides, speaks notes via TTS, detects and "
                "answers audience questions, and integrates with Microsoft Teams. "
                "Let me know if there are any final questions."
            ),
        ),
    ]


# ── Text utilities (TF-IDF) ──────────────────────────────────────────────────

_WORD_RE = re.compile(r"[a-zA-Z0-9]{2,}")

_STOP_WORDS = frozenset(
    "a an the is are was were be been being have has had do does did will would "
    "shall should may might can could of in to for on with at by from and or but "
    "not no nor so yet it its this that these those i we you he she they me us "
    "him her them my our your his".split()
)


def _tokenize(text: str) -> List[str]:
    return [m.group(0).lower() for m in _WORD_RE.finditer(text or "") if m.group(0).lower() not in _STOP_WORDS]


def _compute_idf(documents: List[str]) -> dict[str, float]:
    """Inverse document frequency across the slide corpus."""
    n = len(documents)
    if n == 0:
        return {}
    df: Counter[str] = Counter()
    for doc in documents:
        unique = set(_tokenize(doc))
        for token in unique:
            df[token] += 1
    return {token: math.log((n + 1) / (count + 1)) + 1 for token, count in df.items()}


def _tfidf_similarity(query_tokens: List[str], doc_tokens: List[str], idf: dict[str, float]) -> float:
    """Cosine-ish similarity between query and document using TF-IDF weights."""
    if not query_tokens or not doc_tokens:
        return 0.0

    q_tf = Counter(query_tokens)
    d_tf = Counter(doc_tokens)

    dot = 0.0
    q_norm = 0.0
    for token, qf in q_tf.items():
        w_q = qf * idf.get(token, 1.0)
        w_d = d_tf.get(token, 0) * idf.get(token, 1.0)
        dot += w_q * w_d
        q_norm += w_q ** 2

    d_norm = sum((tf * idf.get(t, 1.0)) ** 2 for t, tf in d_tf.items())

    if q_norm == 0 or d_norm == 0:
        return 0.0
    return dot / (math.sqrt(q_norm) * math.sqrt(d_norm))
